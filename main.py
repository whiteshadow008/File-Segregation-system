import os
import shutil
import zipfile
import joblib
import re
from flask import Flask, render_template, request, redirect, url_for, flash, jsonify
from werkzeug.utils import secure_filename
import PyPDF2
import docx
import openpyxl
from pptx import Presentation # Added missing import

# === Flask Config ===
app = Flask(__name__)
app.secret_key = "supersecretkey"
UPLOAD_FOLDER = "uploads"
SEGREGATED_FOLDER = "segregated"
from tensorflow.keras.models import load_model

MODEL_FILE = "subject_classifier.h5"
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB max file size

# === Load Model ===
if os.path.exists(MODEL_FILE):
    model = load_model(MODEL_FILE)
    print("✅ Deep Learning Model loaded successfully!")
else:
    model = None
    print("⚠️ Model not found! Please train the model first.")

# === Departments & Subjects ===
departments = {
    "AI_DS": ["AI", "ML", "DBMS", "OS", "CN", "Programming", "DSA", "Cloud", "DV"],
    "ECE": [
        "Digital_Electronics", "VLSI", "Signals", "Microprocessors",
        "Communication", "Electronic_Devices", "Embedded", "Control_Systems",
        "EMT", "Optical_Comm", "Antennas", "DSP", "Power_Electronics",
        "Circuit_Theory"
    ],
    "MECH": [
        "Thermodynamics", "CAD", "Mechanics", "Fluids", "Material_Science",
        "Automobile", "Refrigeration", "Manufacturing", "Mechatronics",
        "Kinematics", "Design", "IC_Engines", "FEM", "Robotics", "Vibrations"
    ]
}

# === Ensure Directories ===
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(SEGREGATED_FOLDER, exist_ok=True)
os.makedirs("models", exist_ok=True)

class FileContentExtractor:
    """Extract content from various file types"""
    
    @staticmethod
    def extract_text_from_file(file_path):
        """Extract text content from various file types"""
        try:
            file_ext = os.path.splitext(file_path)[1].lower()
            
            if file_ext == '.pdf':
                return FileContentExtractor._extract_pdf_text(file_path)
            elif file_ext in ['.doc', '.docx']:
                return FileContentExtractor._extract_docx_text(file_path)
            elif file_ext in ['.xls', '.xlsx']:
                return FileContentExtractor._extract_excel_text(file_path)
            elif file_ext in ['.ppt', '.pptx']:
                return FileContentExtractor._extract_pptx_text(file_path)
            elif file_ext == '.txt':
                return FileContentExtractor._extract_txt_text(file_path)
            else:
                # For unsupported files, return filename as content
                return os.path.basename(file_path)
        except Exception as e:
            print(f"Error extracting text from {file_path}: {e}")
            return os.path.basename(file_path)
    
    @staticmethod
    def _extract_pdf_text(file_path):
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text()
        except Exception as e:
            print(f"Error reading PDF: {e}")
        return text
    
    @staticmethod
    def _extract_docx_text(file_path):
        try:
            doc = docx.Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            print(f"Error reading DOCX: {e}")
            return ""
    
    @staticmethod
    def _extract_excel_text(file_path):
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text = ""
            for sheet in workbook.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value:
                            text += str(cell.value) + " "
            return text
        except Exception as e:
            print(f"Error reading Excel: {e}")
            return ""
    
    @staticmethod
    def _extract_pptx_text(file_path):
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            print(f"Error reading PowerPoint: {e}")
            return ""
    
    @staticmethod
    def _extract_txt_text(file_path):
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except:
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    return file.read()
            except Exception as e:
                print(f"Error reading text file: {e}")
                return ""

def clean_text(text):
    """Clean and preprocess text"""
    text = re.sub(r'[^a-zA-Z0-9\s]', ' ', text)
    text = re.sub(r'\s+', ' ', text)
    text = text.lower().strip()
    return text

def classify_and_move(file_path, chosen_subjects):
    """
    Classify file based on content → move into subject folder.
    Returns filename, predicted subject, confidence score, destination folder.
    """
    try:
        filename = os.path.basename(file_path)
        filename_lower = filename.lower()
        
        # New: Check for subject name in filename first
        pre_classified_subject = None
        all_subjects = [subject for sublist in departments.values() for subject in sublist]
        for subject in all_subjects:
            if subject.lower().replace('_', ' ') in filename_lower.replace('_', ' '):
                pre_classified_subject = subject
                break
        
        if pre_classified_subject:
            prediction = pre_classified_subject
            confidence = 100.0
            print(f"File: {filename}")
            print(f"Pre-classified by filename: {prediction} (Confidence: {confidence}%)")
        else:
            # Original logic: Extract content and predict with model
            content = FileContentExtractor.extract_text_from_file(file_path)
            
            if not content.strip():
                content = filename
            
            cleaned_content = clean_text(content)
            
            if model and cleaned_content.strip():
                prediction = model.predict([cleaned_content])[0]
                probabilities = model.predict_proba([cleaned_content])[0]
                confidence = round(max(probabilities) * 100, 2)
                
                print(f"File: {filename}")
                print(f"Content preview: {content[:100]}...")
                print(f"Predicted by model: {prediction} (Confidence: {confidence}%)")
            else:
                prediction, confidence = "Uncategorized", 0.0
                print(f"Model not available or no content found for: {filename}")

        # If prediction not in chosen subjects → put in Corrected
        if chosen_subjects and prediction not in chosen_subjects:
            folder = "Corrected"
        else:
            folder = prediction

        # Create destination directory and move file
        dest_dir = os.path.join(SEGREGATED_FOLDER, folder)
        os.makedirs(dest_dir, exist_ok=True)
        
        # Copy file to destination
        dest_path = os.path.join(dest_dir, filename)
        shutil.copy2(file_path, dest_path)

        return filename, prediction, confidence, folder
    
    except Exception as e:
        print(f"Error processing file {file_path}: {e}")
        # Move to uncategorized on error
        folder = "Uncategorized"
        dest_dir = os.path.join(SEGREGATED_FOLDER, folder)
        os.makedirs(dest_dir, exist_ok=True)
        shutil.copy2(file_path, dest_dir)
        return filename, "Error", 0.0, folder

@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        department = request.form.get("department")
        chosen_subjects = request.form.getlist("subjects")
        upload_type = request.form.get("upload_type")

        if not department or not chosen_subjects:
            flash("Please select a department and subjects!", "error")
            return redirect(url_for("index"))

        # Clear old files
        if os.path.exists(UPLOAD_FOLDER):
            shutil.rmtree(UPLOAD_FOLDER)
        os.makedirs(UPLOAD_FOLDER, exist_ok=True)

        if os.path.exists(SEGREGATED_FOLDER):
            shutil.rmtree(SEGREGATED_FOLDER)
        os.makedirs(SEGREGATED_FOLDER, exist_ok=True)

        uploaded_files = []
        
        try:
            if upload_type == "file":
                files = request.files.getlist("file")
                for file in files:
                    if file and file.filename:
                        filename = secure_filename(file.filename)
                        file_path = os.path.join(UPLOAD_FOLDER, filename)
                        file.save(file_path)
                        uploaded_files.append(file_path)

            elif upload_type == "folder":
                files = request.files.getlist("folder")
                for file in files:
                    if file and file.filename:
                        filename = secure_filename(file.filename)
                        file_path = os.path.join(UPLOAD_FOLDER, filename)
                        file.save(file_path)
                        uploaded_files.append(file_path)

            elif upload_type == "zip":
                zip_file = request.files.get("zip_file")
                if zip_file and zip_file.filename.endswith('.zip'):
                    zip_path = os.path.join(UPLOAD_FOLDER, secure_filename(zip_file.filename))
                    zip_file.save(zip_path)
                    
                    # Extract zip file
                    with zipfile.ZipFile(zip_path, 'r') as zip_ref:
                        zip_ref.extractall(UPLOAD_FOLDER)
                    
                    # Get all extracted files
                    for root, dirs, files in os.walk(UPLOAD_FOLDER):
                        for file in files:
                            if not file.endswith('.zip'):
                                uploaded_files.append(os.path.join(root, file))

            if not uploaded_files:
                flash("No files were uploaded!", "error")
                return redirect(url_for("index"))

            # Run classification
            results = []
            summary = {}
            
            for file_path in uploaded_files:
                try:
                    fname, pred, conf, folder = classify_and_move(file_path, chosen_subjects)
                    results.append((fname, pred, conf, folder))
                    summary[folder] = summary.get(folder, 0) + 1
                except Exception as e:
                    print(f"Error processing {file_path}: {e}")
                    results.append((os.path.basename(file_path), "Error", 0.0, "Uncategorized"))
                    summary["Uncategorized"] = summary.get("Uncategorized", 0) + 1

            if results:
                flash(f"Successfully processed {len(results)} files!", "success")
                return render_template("result.html", results=results, summary=summary,
                             department=department, chosen_subjects=chosen_subjects)

            else:
                flash("No files were processed!", "error")
                return redirect(url_for("index"))
                
        except Exception as e:
            flash(f"Error processing files: {str(e)}", "error")
            return redirect(url_for("index"))

    return render_template("index.html", departments=departments)

@app.route("/predict", methods=["POST"])
def predict_text():
    """API endpoint to predict subject from text"""
    data = request.json
    text = data.get('text', '')
    
    if not text.strip():
        return jsonify({'error': 'No text provided'}), 400
    
    if not model:
        return jsonify({'error': 'Model not loaded'}), 500
    
    try:
        cleaned_text = clean_text(text)
        prediction = model.predict([cleaned_text])[0]
        probabilities = model.predict_proba([cleaned_text])[0]
        confidence = round(max(probabilities) * 100, 2)
        
        return jsonify({
            'prediction': prediction,
            'confidence': confidence,
            'success': True
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route("/model_status")
def model_status():
    """Check if model is loaded"""
    return jsonify({
        'model_loaded': model is not None,
        'model_classes': list(model.classes_) if model else []
    })

if __name__ == "__main__":
    app.run(debug=True, port=5000)
