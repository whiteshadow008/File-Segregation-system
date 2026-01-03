import os
import shutil
import zipfile
import pickle
from flask import Flask, render_template, request, redirect, url_for, flash
from werkzeug.utils import secure_filename

# File extraction
import PyPDF2
import docx
import openpyxl
from pptx import Presentation
import re
import numpy as np
from tensorflow.keras.models import load_model
from tensorflow.keras.preprocessing.sequence import pad_sequences

# -------------------------------
# Flask Config
# -------------------------------
app = Flask(__name__)
app.secret_key = "supersecretkey"
UPLOAD_FOLDER = "uploads"
SEGREGATED_FOLDER = "segregated"
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(SEGREGATED_FOLDER, exist_ok=True)

# -------------------------------
# Load RNN Model, Tokenizer, Encoder, max_len
# -------------------------------
MODEL_FILE = "rnn_text_classifier.h5"
TOKENIZER_FILE = "tokenizer.pkl"
ENCODER_FILE = "label_encoder.pkl"
MAX_LEN_FILE = "max_len.pkl"

try:
    model = load_model(MODEL_FILE)
    with open(TOKENIZER_FILE, 'rb') as f:
        tokenizer = pickle.load(f)
    with open(ENCODER_FILE, 'rb') as f:
        encoder = pickle.load(f)
    with open(MAX_LEN_FILE, 'rb') as f:
        max_len = pickle.load(f)
    print("✅ Model, Tokenizer, Encoder, max_len loaded successfully!")
except Exception as e:
    print("⚠️ Failed to load model or tokenizer/encoder/max_len:", e)
    model, tokenizer, encoder, max_len = None, None, None, None

# -------------------------------
# Departments & Subjects
# -------------------------------
departments = {
    "AI_DS": ["AI", "ML", "DBMS", "OS", "CN", "Programming", "DSA", "Cloud", "DV"],
    "ECE": ["Digital_Electronics", "VLSI", "Signals", "Microprocessors",
            "Communication", "Electronic_Devices", "Embedded", "Control_Systems",
            "EMT", "Optical_Comm", "Antennas", "DSP", "Power_Electronics",
            "Circuit_Theory"],
    "MECH": ["Thermodynamics", "CAD", "Mechanics", "Fluids", "Material_Science",
             "Automobile", "Refrigeration", "Manufacturing", "Mechatronics",
             "Kinematics", "Design", "IC_Engines", "FEM", "Robotics", "Vibrations"]
}

# -------------------------------
# File Content Extractor
# -------------------------------
class FileContentExtractor:
    @staticmethod
    def extract_text_from_file(file_path):
        ext = os.path.splitext(file_path)[1].lower()
        try:
            if ext == ".pdf":
                text = ""
                with open(file_path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        text += page.extract_text() or ""
                return text
            elif ext in [".doc", ".docx"]:
                doc = docx.Document(file_path)
                return "\n".join([p.text for p in doc.paragraphs])
            elif ext in [".xls", ".xlsx"]:
                wb = openpyxl.load_workbook(file_path, data_only=True)
                text = ""
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows():
                        for cell in row:
                            if cell.value:
                                text += str(cell.value) + " "
                return text
            elif ext in [".ppt", ".pptx"]:
                prs = Presentation(file_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            elif ext == ".txt":
                try:
                    with open(file_path, "r", encoding="utf-8") as f:
                        return f.read()
                except:
                    with open(file_path, "r", encoding="latin-1") as f:
                        return f.read()
            else:
                return os.path.basename(file_path)
        except Exception as e:
            print(f"Error extracting text from {file_path}: {e}")
            return os.path.basename(file_path)

# -------------------------------
# Text Cleaning
# -------------------------------
def clean_text(text):
    text = re.sub(r'[^a-zA-Z0-9\s]', ' ', text)
    text = re.sub(r'\s+', ' ', text).lower().strip()
    return text

# -------------------------------
# Predict & Move File
# -------------------------------
def classify_and_move(file_path, chosen_subjects):
    filename = os.path.basename(file_path)
    content = FileContentExtractor.extract_text_from_file(file_path)
    
    # Combine filename + content
    combined_text = clean_text(filename + " " + content)
    
    confidence = 0.0
    predicted_subject = "Uncategorized"

    # First, check if filename contains a subject hint
    all_subjects = [s for sublist in departments.values() for s in sublist]
    for subj in all_subjects:
        if subj.lower().replace('_', ' ') in filename.lower().replace('_', ' '):
            predicted_subject = subj
            confidence = 100.0
            break
    else:
        # If no filename match, use RNN
        if model and combined_text.strip():
            seq = tokenizer.texts_to_sequences([combined_text])
            padded = pad_sequences(seq, maxlen=max_len, padding='post')
            pred_probs = model.predict(padded, verbose=0)
            class_index = np.argmax(pred_probs, axis=1)
            predicted_subject = encoder.inverse_transform(class_index)[0]
            confidence = round(float(np.max(pred_probs) * 100), 2)

    # If predicted_subject not in chosen subjects → move to Corrected
    folder = predicted_subject if (predicted_subject in chosen_subjects) else "Corrected"
    dest_dir = os.path.join(SEGREGATED_FOLDER, folder)
    os.makedirs(dest_dir, exist_ok=True)
    shutil.copy2(file_path, os.path.join(dest_dir, filename))
    
    return filename, predicted_subject, confidence, folder

# -------------------------------
# Routes
# -------------------------------
@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        department = request.form.get("department")
        chosen_subjects = request.form.getlist("subjects")
        upload_type = request.form.get("upload_type")

        if not department or not chosen_subjects:
            flash("Please select department and subjects!", "error")
            return redirect(url_for("index"))

        # Clear uploads & segregated folders
        shutil.rmtree(UPLOAD_FOLDER, ignore_errors=True)
        shutil.rmtree(SEGREGATED_FOLDER, ignore_errors=True)
        os.makedirs(UPLOAD_FOLDER, exist_ok=True)
        os.makedirs(SEGREGATED_FOLDER, exist_ok=True)

        uploaded_files = []

        # File upload handling
        if upload_type == "file":
            files = request.files.getlist("file")
            for file in files:
                if file and file.filename:
                    fname = secure_filename(file.filename)
                    path = os.path.join(UPLOAD_FOLDER, fname)
                    file.save(path)
                    uploaded_files.append(path)
        elif upload_type == "folder":
                files = request.files.getlist("folder")
                for file in files:
                    if file and file.filename:
                        filename = secure_filename(file.filename)
                        file_path = os.path.join(UPLOAD_FOLDER, filename)
                        file.save(file_path)
                        uploaded_files.append(file_path)
        elif upload_type == "zip":
            zip_file = request.files.get("zip_file")
            if zip_file and zip_file.filename.endswith('.zip'):
                zip_path = os.path.join(UPLOAD_FOLDER, secure_filename(zip_file.filename))
                zip_file.save(zip_path)
                with zipfile.ZipFile(zip_path, 'r') as zip_ref:
                    zip_ref.extractall(UPLOAD_FOLDER)
                for root, dirs, files in os.walk(UPLOAD_FOLDER):
                    for file in files:
                        if not file.endswith('.zip'):
                            uploaded_files.append(os.path.join(root, file))

        if not uploaded_files:
            flash("No files uploaded!", "error")
            return redirect(url_for("index"))

        # Classify all files
        results = []
        summary = {}
        for file_path in uploaded_files:
            fname, pred, conf, folder = classify_and_move(file_path, chosen_subjects)
            results.append((fname, pred, conf, folder))
            summary[folder] = summary.get(folder, 0) + 1

        return render_template("result.html", results=results, summary=summary,
                               department=department, chosen_subjects=chosen_subjects)

    return render_template("index.html", departments=departments)

# -------------------------------
# Run Flask
# -------------------------------
if __name__ == "__main__":
    app.run(debug=True, port=8000)
